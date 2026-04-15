"""
Generate a PowerPoint slide presentation from script scenes.
Supports image-only slides (full-screen) or text-only fallback.
Usage: python -m execution.slide_generator --scenes_json <path> --output_path <path>
"""
import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from execution.config import TMP_DIR


BG_COLOR = RGBColor(0x0A, 0x0A, 0x0A)       # near-black
TITLE_COLOR = RGBColor(0xFF, 0xD7, 0x00)     # gold
TEXT_COLOR = RGBColor(0xF0, 0xF0, 0xF0)      # off-white
ACCENT_COLOR = RGBColor(0x7C, 0x3A, 0xED)    # purple
SCENE_NUM_COLOR = RGBColor(0xA0, 0xA0, 0xA0) # gray


def _set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def generate_slides(scenes: list[dict], output_path: str = None,
                    title: str = "Script Presentation",
                    scene_images: list[dict] = None) -> str:
    """
    Create a PPTX presentation from scenes.
    If scene_images provided: full-screen image slides (no text).
    If no images: text-only slides with dark theme.
    """
    output_path = output_path or str(TMP_DIR / "script_slides.pptx")
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    # Build image lookup
    img_lookup = {}
    if scene_images:
        for si in scene_images:
            if si.get("image_path") and Path(si["image_path"]).exists():
                img_lookup[si["scene"]] = si["image_path"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ── Title slide ──
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(title_slide, BG_COLOR)

    txBox = title_slide.shapes.add_textbox(
        Inches(1), Inches(2.5), slide_width - Inches(2), Inches(2)
    )
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    sub_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(4.8), slide_width - Inches(2), Inches(1)
    )
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = f"{len(scenes)} Scenes"
    sp.font.size = Pt(24)
    sp.font.color.rgb = SCENE_NUM_COLOR
    sp.alignment = PP_ALIGN.CENTER

    # ── Scene slides ──
    for scene in scenes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        scene_num = scene.get("scene_number", "")
        scene_text = scene.get("text", "")
        image_path = img_lookup.get(scene_num)

        if image_path:
            # ══════ IMAGE-ONLY SLIDE ══════
            # Full-screen image, no text
            _set_slide_bg(slide, BG_COLOR)

            try:
                # Add image stretched to full slide
                slide.shapes.add_picture(
                    image_path,
                    Emu(0), Emu(0),
                    width=slide_width,
                    height=slide_height
                )
            except Exception:
                # If image fails, fall back to text
                _add_text_slide_content(slide, scene_num, scene_text,
                                        len(scenes), slide_width, slide_height)

            # Small scene number overlay (bottom-right, semi-transparent feel)
            pg_box = slide.shapes.add_textbox(
                slide_width - Inches(1.2), slide_height - Inches(0.5),
                Inches(1), Inches(0.4)
            )
            pgp = pg_box.text_frame.paragraphs[0]
            pgp.text = f"{scene_num}/{len(scenes)}"
            pgp.font.size = Pt(11)
            pgp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            pgp.font.bold = True
            pgp.alignment = PP_ALIGN.RIGHT
        else:
            # ══════ TEXT-ONLY SLIDE (fallback) ══════
            _set_slide_bg(slide, BG_COLOR)
            _add_text_slide_content(slide, scene_num, scene_text,
                                    len(scenes), slide_width, slide_height)

    # ── End slide ──
    end_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(end_slide, BG_COLOR)
    end_box = end_slide.shapes.add_textbox(
        Inches(1), Inches(2.5), slide_width - Inches(2), Inches(2)
    )
    ep = end_box.text_frame.paragraphs[0]
    ep.text = "End"
    ep.font.size = Pt(48)
    ep.font.bold = True
    ep.font.color.rgb = TITLE_COLOR
    ep.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path


def _add_text_slide_content(slide, scene_num, scene_text, total_scenes,
                            slide_width, slide_height):
    """Add text-only content to a slide (used when no image available)."""
    # Accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.15), slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # Scene number
    num_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4), Inches(3), Inches(0.5)
    )
    np = num_box.text_frame.paragraphs[0]
    np.text = f"SCENE {scene_num}"
    np.font.size = Pt(14)
    np.font.bold = True
    np.font.color.rgb = SCENE_NUM_COLOR

    # Scene text
    text_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.2),
        slide_width - Inches(1.5), slide_height - Inches(2)
    )
    ttf = text_box.text_frame
    ttf.word_wrap = True

    for i, para_text in enumerate(scene_text.split("\n")):
        p = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
        p.text = para_text.strip()
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)

    # Page number
    pg_box = slide.shapes.add_textbox(
        slide_width - Inches(1.5), slide_height - Inches(0.6),
        Inches(1.2), Inches(0.4)
    )
    pgp = pg_box.text_frame.paragraphs[0]
    pgp.text = f"{scene_num} / {total_scenes}"
    pgp.font.size = Pt(12)
    pgp.font.color.rgb = SCENE_NUM_COLOR
    pgp.alignment = PP_ALIGN.RIGHT


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate slide presentation")
    parser.add_argument("--scenes_json", required=True)
    parser.add_argument("--images_json", default=None)
    parser.add_argument("--output_path", default=str(TMP_DIR / "script_slides.pptx"))
    parser.add_argument("--title", default="Script Presentation")
    args = parser.parse_args()

    with open(args.scenes_json, "r", encoding="utf-8") as f:
        scenes = json.load(f)

    images = None
    if args.images_json:
        with open(args.images_json, "r", encoding="utf-8") as f:
            images = json.load(f)

    result = generate_slides(scenes, args.output_path, args.title, images)
    print(f"Slides saved to: {result}")
